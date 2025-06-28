import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os

from pptx import Presentation
from PIL import Image
import pytesseract
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain_community.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv

load_dotenv()
api_key = os.getenv("GOOGLE_API_KEY")
if not api_key:
    st.error("GOOGLE_API_KEY not found in environment variables!")
    st.stop()

genai.configure(api_key=api_key)

pytesseract.pytesseract.tesseract_cmd = '/opt/homebrew/bin/tesseract'

def get_text(file):
    text = ""
    file_extension = os.path.splitext(file.name)[1].lower()
    try:
        if file_extension == '.pdf':
            pdf_reader = PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text()
                
        elif file_extension == '.pptx':
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                        
        elif file_extension in ('.txt', '.md'):
            text = file.read().decode("utf-8")
            
        elif file_extension in ('.jpg', '.jpeg', '.png'):
            image = Image.open(file)
            text = pytesseract.image_to_string(image)
            
        else:
            st.warning(f"Unsupported file format: {file_extension}")
            
    except Exception as e:
        st.error(f"Error processing {file.name}: {str(e)}")
        
    return text

def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

def get_files_text(files):
    """Process multiple files and extract text"""
    text = ""
    for file in files:
        text += get_text(file) + "\n\n"
    return text

def get_vector_store(text_chunks):
    try:
        embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
        vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
        vector_store.save_local("faiss_index")
        return True
    except Exception as e:
        st.error(f"Error creating vector store: {str(e)}")
        return False

def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context, make sure to provide all the details, if the answer is not in
    provided context just say, "answer is not available in the context", don't provide the wrong answer\n\n
    Context:\n {context}?\n
    Question: \n{question}\n

    Answer:
    """

    try:
        # Try different model names in order of preference
        model_names = [
            "gemini-1.5-flash",
            "gemini-1.5-pro", 
            "gemini-1.0-pro",
            "models/gemini-1.5-flash",
            "models/gemini-1.5-pro",
            "models/gemini-1.0-pro"
        ]
        
        model = None
        for model_name in model_names:
            try:
                model = ChatGoogleGenerativeAI(model=model_name, temperature=0.3)
                # Test the model with a simple call
                break
            except Exception as e:
                continue
        
        if model is None:
            st.error("Could not initialize any Gemini model. Please check your API key and model availability.")
            return None
            
        prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
        chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
        return chain
        
    except Exception as e:
        st.error(f"Error creating conversational chain: {str(e)}")
        return None

def user_input(user_question):
    try:
        embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
        
        # Check if faiss_index exists
        if not os.path.exists("faiss_index"):
            st.error("Please upload and process PDF files first!")
            return
            
        new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
        docs = new_db.similarity_search(user_question)

        chain = get_conversational_chain()
        if chain is None:
            return

        response = chain(
            {"input_documents": docs, "question": user_question},
            return_only_outputs=True
        )

        print(response)
        st.write("Reply: ", response["output_text"])
        
    except Exception as e:
        st.error(f"Error processing question: {str(e)}")

def list_available_models():
    """Helper function to list available models"""
    try:
        models = []
        for model in genai.list_models():
            if 'generateContent' in model.supported_generation_methods:
                models.append(model.name)
        return models
    except Exception as e:
        st.error(f"Error listing models: {str(e)}")
        return []

def main():
    st.set_page_config(page_title="DocDecoder", layout="wide")
    st.header("DocDecoder")

    user_question = st.text_input("Ask a Question from the Files")

    if user_question:
        user_input(user_question)

    with st.sidebar:
        st.title("Menu:")
        file = st.file_uploader(
            "Upload your Files and Click on the Submit & Process Button", 
            accept_multiple_files=True,
            type=['pdf', 'pptx', 'txt', 'jpg', 'jpeg', 'png', 'md']
        )
        
        if st.button("Submit & Process"):
            if not file:
                st.error("Please upload at least one file!")
            else:
                with st.spinner("Processing..."):
                    try:
                        raw_text = get_files_text(file)
                        if not raw_text.strip():
                            st.error("No text could be extracted from the files!")
                            return
                            
                        text_chunks = get_text_chunks(raw_text)
                        if get_vector_store(text_chunks):
                            st.success("Files processed successfully! You can now ask questions.")
                        else:
                            st.error("Failed to process files!")
                    except Exception as e:
                        st.error(f"Error processing files: {str(e)}")

if __name__ == "__main__":
    main()
